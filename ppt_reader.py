import streamlit as st
from pptx import Presentation
from gtts import gTTS
import tempfile
import os

# 음성 선택 옵션 (남성, 여성)
LANGUAGE = 'ko'  # 한국어 설정

# PPT 파일 업로드
st.title('PPT 대본 수정 및 음성 변환 웹앱')

uploaded_file = st.file_uploader("PPT 파일 업로드", type="pptx")
if uploaded_file is not None:
    presentation = Presentation(uploaded_file)

    # 각 슬라이드의 텍스트 추출 및 수정
    slide_scripts = {}
    for i, slide in enumerate(presentation.slides):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
        # 대본 수정할 수 있도록 텍스트 박스 생성
        slide_scripts[i] = st.text_area(f"슬라이드 {i+1} 대본 수정", text)
    
    # 음성 성별 선택 (gTTS는 성별을 선택하지 않으므로 한국어 음성은 고정)
    # 한국어에서 남성, 여성 음성 차이는 제공되지 않음
    st.write("현재 한국어 음성은 남성과 여성 구분 없이 제공됩니다.")
    
    # 음성 파일 생성
    if st.button("음성 파일 생성"):
        for i, script in slide_scripts.items():
            # gTTS로 텍스트를 음성으로 변환
            tts = gTTS(text=script, lang=LANGUAGE)
            
            # 임시 파일에 음성 저장
            temp_audio_path = f"/tmp/slide_{i+1}.mp3"
            tts.save(temp_audio_path)
            
            # 생성된 음성 파일 재생
            st.audio(temp_audio_path)
        
        st.success("음성 파일이 생성되었습니다.")
